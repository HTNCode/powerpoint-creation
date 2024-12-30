from gpt_researcher import GPTResearcher
import asyncio
import os
from typing import Dict, Any, List, Tuple, Optional, Protocol, Union, cast
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Run, _Paragraph
from pptx.shapes.placeholder import SlidePlaceholder
from openai import AsyncOpenAI
import aiohttp
from yarl import URL
from pptx.shapes.placeholder import SlidePlaceholder

class TextFrameProtocol(Protocol):
    def clear(self) -> None: ...
    def add_paragraph(self) -> _Paragraph: ...
    @property
    def paragraphs(self) -> List[_Paragraph]: ...

class FontProtocol(Protocol):
    @property
    def name(self) -> str: ...
    @name.setter
    def name(self, value: str) -> None: ...
    @property
    def size(self) -> Length: ...
    @size.setter
    def size(self, value: Length) -> None: ...

class ParagraphProtocol(Protocol):
    @property
    def text(self) -> str: ...
    @text.setter
    def text(self, value: str) -> None: ...
    @property
    def font(self) -> FontProtocol: ...
    @property
    def level(self) -> int: ...
    @level.setter
    def level(self, value: int) -> None: ...

class ShapeProtocol(Protocol):
    @property
    def text_frame(self) -> TextFrameProtocol: ...
    @property
    def text(self) -> str: ...

class PlaceholderProtocol(Protocol):
    @property
    def text_frame(self) -> TextFrameProtocol: ...
    @property
    def text(self) -> str: ...


def safe_get_text_frame(
    shape: Optional[Union[ShapeProtocol, Shape, PlaceholderProtocol, SlidePlaceholder]]
) -> Optional[TextFrameProtocol]:
    """安全にtext_frameを取得するヘルパー関数"""
    if shape is None:
        return None
    if not hasattr(shape, 'text_frame'):
        return None
    return cast(TextFrameProtocol, shape.text_frame)

class PowerPointGenerator:
    def __init__(self, api_key: str):
        """PowerPointGeneratorの初期化"""
        self.prs = Presentation()
        self._setup_slide_layouts()
        self.client = AsyncOpenAI(api_key=api_key)

    def _setup_slide_layouts(self):
        """スライドのレイアウトを設定"""
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        # スライドマスターの設定を調整
        for layout in self.prs.slide_layouts:
            for placeholder in layout.placeholders:
                try:
                    text_frame = safe_get_text_frame(placeholder)
                    if text_frame:
                        for paragraph in text_frame.paragraphs:
                            p = cast(ParagraphProtocol, paragraph)
                            p.font.name = 'BIZ UDPゴシック'
                except Exception:
                    continue

    async def _generate_and_save_image(self, prompt: str, image_path: str) -> Optional[str]:
        """DALL-E 3を使用して画像を生成して保存"""
        try:
            response = await self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size="1792x1024",
                quality="standard",
                n=1,
            )

            image_url = str(response.data[0].url)
            if not image_url:
                return None

            async with aiohttp.ClientSession() as session:
                async with session.get(URL(image_url)) as response:
                    if response.status == 200:
                        with open(image_path, 'wb') as f:
                            f.write(await response.read())
                        return image_path
                    else:
                        print(f"画像のダウンロードに失敗しました: {response.status}")
                        return None
        except Exception as e:
            print(f"画像生成エラー: {str(e)}")
            return None

    def _add_title_slide(self, title: str, output_dir: str):
        """タイトルスライドを追加"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        # タイトルの設定
        text_frame = safe_get_text_frame(slide.shapes.title)
        if text_frame:
            try:
                text_frame.clear()
                p = cast(ParagraphProtocol, text_frame.paragraphs[0])
                p.text = title
                p.font.name = 'BIZ UDPゴシック'
                p.font.size = Pt(32)
            except Exception:
                print("タイトルの設定に失敗しました")

        # 背景画像の設定
        image_path = os.path.join(output_dir, "title_image.png")
        if os.path.exists(image_path):
            try:
                left = Inches(0)
                top = Inches(0)
                width = self.prs.slide_width
                height = self.prs.slide_height
                picture = slide.shapes.add_picture(image_path, left, top, width, height)
                try:
                    # 画像を最背面に移動
                    slide.shapes._spTree.insert(0, picture._element)
                except Exception:
                    print("画像の重ね順の設定に失敗しました")
            except Exception as e:
                print(f"背景画像の設定に失敗しました: {str(e)}")

    async def _add_content_slide(self, title: str, content: str, output_dir: str, slide_num: int):
        """コンテンツスライドを追加"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        # タイトルの設定
        text_frame = safe_get_text_frame(slide.shapes.title)
        if text_frame:
            try:
                text_frame.clear()
                p = cast(ParagraphProtocol, text_frame.paragraphs[0])
                p.text = title
                p.font.name = 'BIZ UDPゴシック'
                p.font.size = Pt(28)
            except Exception:
                print("スライドタイトルの設定に失敗しました")

        # 本文の設定
        if len(slide.placeholders) > 1:
            placeholder = cast(PlaceholderProtocol, slide.placeholders[1])
            text_frame = safe_get_text_frame(placeholder)
            if text_frame:
                try:
                    text_frame.clear()
                    for para_text in content.split('\n'):
                        if para_text.strip():
                            p = cast(ParagraphProtocol, text_frame.add_paragraph())
                            p.text = para_text
                            p.font.name = 'BIZ UDPゴシック'
                            p.font.size = Pt(16)
                            if para_text.startswith('•') or para_text.startswith('-'):
                                p.level = 1
                            else:
                                p.level = 0
                except Exception:
                    print("本文の設定に失敗しました")

    def _parse_research_content(self, content: str) -> Tuple[str, List[Dict[str, str]]]:
        """研究内容をスライド用に解析"""
        slides_data = []

        # Marpのフロントマター（---で囲まれた部分）を除去
        content_parts = content.split('---\n')
        if len(content_parts) > 2:
            # フロントマターを除いた部分を使用
            content = '---\n'.join(content_parts[2:])
        else:
            content = content_parts[-1]

        # スライドを正しく分割
        slides = content.split('\n---\n')

        # 最初のスライド（タイトルスライド）の処理
        first_slide = slides[0].strip()
        lines = first_slide.split('\n')
        title = ''
        content_lines = []

        for line in lines:
            if line.startswith('# '):
                title = line.replace('# ', '').strip()
            elif line.strip():
                content_lines.append(line.strip())

        if not title:
            title = "無題のプレゼンテーション"

        if content_lines:
            slides_data.append({
                'title': 'はじめに',
                'content': '\n'.join(content_lines)
            })

        # 残りのスライドの処理
        for slide in slides[1:]:
            if not slide.strip():
                continue

            lines = slide.strip().split('\n')
            slide_title = ''
            slide_content_lines = []

            for line in lines:
                if line.startswith('## '):
                    slide_title = line.replace('## ', '').strip()
                elif line.strip():
                    # 箇条書きを日本語スタイルに調整
                    line = line.replace('* ', '• ').replace('- ', '• ')
                    slide_content_lines.append(line.strip())

            if slide_title and slide_content_lines:
                slides_data.append({
                    'title': slide_title,
                    'content': '\n'.join(slide_content_lines)
                })
            elif slide_content_lines:  # タイトルがない場合
                slides_data.append({
                    'title': '続き',
                    'content': '\n'.join(slide_content_lines)
                })

        return title, slides_data

    async def create_presentation(self, markdown_content: str, output_dir: str) -> str:
        """PowerPointプレゼンテーションを作成"""
        os.makedirs(output_dir, exist_ok=True)

        try:
            title, slides_data = self._parse_research_content(markdown_content)

            # タイトルスライド用の画像生成プロンプト
            title_image_prompt = f"""
            以下のトピックに関するプレゼンテーション表紙の画像を作成:
            {title}

            スタイル:
            - モダンでプロフェッショナルなデザイン
            - ビジネスプレゼンテーションに適した抽象的な背景
            - 清潔で洗練された印象
            - 日本のビジネス文化に適した控えめな配色
            """

            await self._generate_and_save_image(
                title_image_prompt,
                os.path.join(output_dir, "title_image.png")
            )

            self._add_title_slide(title, output_dir)

            for i, slide_data in enumerate(slides_data):
                await self._add_content_slide(
                    slide_data['title'],
                    slide_data['content'],
                    output_dir,
                    i + 1
                )

            pptx_path = os.path.join(output_dir, "presentation.pptx")
            self.prs.save(pptx_path)

            # ファイルが正しく保存されたことを確認
            if not os.path.exists(pptx_path) or os.path.getsize(pptx_path) == 0:
                raise Exception("プレゼンテーションファイルの保存に失敗しました")

            return pptx_path

        except Exception as e:
            print(f"プレゼンテーション作成中にエラーが発生しました: {str(e)}")
            raise

async def get_report(query: str, report_type: str) -> str:
    researcher = GPTResearcher(
        query=query,
        report_type=report_type,
        report_format="markdown"
    )
    await researcher.conduct_research()
    report = await researcher.write_report()
    return report

async def translate_report(report: str, api_key: str) -> str:
    """レポートを日本語に翻訳する関数"""
    client = AsyncOpenAI(api_key=api_key)
    try:
        response = await client.chat.completions.create(
            model="gpt-4-turbo-preview",
            messages=[
                {
                    "role": "system",
                    "content": """あなたは英語から日本語への翻訳の専門家です。
                        入力された英語の文章を自然な日本語に翻訳してください。
                        Marp形式を維持して出力してください。具体的には：
                        - フロントマター（---で囲まれた設定部分）は保持
                        - スライド区切り（---）は保持
                        - 各スライドの見出し（#や##）は保持
                        - 箇条書きの形式は保持"""
                },
                {
                    "role": "user",
                    "content": report
                }
            ],
            temperature=0.3
        )
        translated_content = response.choices[0].message.content
        if translated_content is None:
            raise ValueError("翻訳結果が空でした")
        print("翻訳が完了しました")
        return translated_content
    except Exception as e:
        print(f"翻訳中にエラーが発生しました: {str(e)}")
        raise

async def main():
    try:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY environment variable is not set")

        query = input("研究するトピックを入力してください: ")
        report_type = "research_report"

        print("研究を開始しています...")
        report = await get_report(query, report_type)

        print("翻訳を実行しています...")
        translated_report = await translate_report(report, api_key)

        print("プレゼンテーションを作成しています...")
        generator = PowerPointGenerator(api_key)
        output_dir = "output"

        pptx_path = await generator.create_presentation(translated_report, output_dir)
        print(f"プレゼンテーションが作成されました: {pptx_path}")

    except Exception as e:
        print(f"エラーが発生しました: {str(e)}")
        raise

if __name__ == "__main__":
    asyncio.run(main())